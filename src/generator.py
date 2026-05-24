"""PPTXジェネレーター - スライド生成コアモジュール"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path


def create_presentation(slides_data: list[dict], output_path: str, design: dict | None = None) -> str:
    """スライドデータからPPTXファイルを生成する。

    Args:
        slides_data: スライド定義のリスト。各要素は以下のキーを持つ:
            - layout: "title", "content", "section", "blank"
            - title: スライドタイトル
            - body: 本文テキストまたはリスト
            - notes: スピーカーノート(任意)
        output_path: 出力ファイルパス
        design: デザイン設定(任意)
            - font_name: フォント名
            - title_size: タイトルフォントサイズ(Pt)
            - body_size: 本文フォントサイズ(Pt)

    Returns:
        生成されたファイルのパス
    """
    prs = Presentation()
    design = design or {}
    font_name = design.get("font_name", "Meiryo")
    title_size = Pt(design.get("title_size", 32))
    body_size = Pt(design.get("body_size", 18))

    layout_map = {
        "title": 0,
        "content": 1,
        "section": 2,
        "blank": 6,
    }

    for slide_data in slides_data:
        layout_idx = layout_map.get(slide_data.get("layout", "content"), 1)
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)

        # タイトル設定
        if slide.shapes.title and slide_data.get("title"):
            tf = slide.shapes.title.text_frame
            tf.text = slide_data["title"]
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.name = font_name
                    run.font.size = title_size

        # 本文設定
        body = slide_data.get("body")
        if body and len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            lines = body if isinstance(body, list) else [body]
            for i, line in enumerate(lines):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.text = line
                para.font.name = font_name
                para.font.size = body_size

        # スピーカーノート
        if slide_data.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
